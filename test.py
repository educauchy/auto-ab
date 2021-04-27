import statsmodels.stats.api as sms
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches


create_presentation = False
save_plots = False
power_test = []

if create_presentation:
    prs=Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

segments = np.array([-1, 0, 1, 2, 4, 5, 7, 8, 9, 10, 16, 17, 19, 27, 30, 33, 47, 48, 50, 51, 55, 59, 68])
stds = np.array([281, 286, 280, 172, 280, 144, 284, 310, 265, 290, 276, 341, 208, 261, 292, 304, 332, 281, 74, 199, 176, 207, 68])
means = np.array([288, 287, 294, 206, 279, 192, 295, 301, 273, 314, 245, 350, 246, 294, 340, 331, 309, 293, 248, 215, 246, 196, 232])
conv = np.array([35349, 4992, 4876, 264, 2619, 625, 1970, 2697, 5581, 1437, 2187, 987, 1368, 740, 653, 1979, 183, 778, 70, 83, 429, 599, 224])

print('Check array sizes:')
print(segments.size == stds.size == means.size == conv.size)

effect_size = np.array([*range(5, 101, 5)])

for index, segment in enumerate(segments):
    segment_name = segment if segment >= 0 else 'all'
    n_samples_arr = []
    for es in effect_size:
        n_samples = sms.TTestIndPower().solve_power(
            effect_size=es / stds[index],
            power=0.9,
            nobs1=None,
            alpha=0.05,
            ratio=1
        )
        n_samples_arr.append(n_samples)

        check_eff_size = sms.TTestIndPower().solve_power(
            effect_size=None,
            power=0.9,
            nobs1=n_samples,
            alpha=0.05,
            ratio=1
        )

        power_test.append(round(es/stds[index], 3) == round(check_eff_size, 3))

    eff_size_to_mean = sms.TTestIndPower().solve_power(
        effect_size=None,
        power=0.9,
        nobs1=conv[index],
        alpha=0.05,
        ratio=1
    )


    if save_plots:
        fig, (ax, ax_table) = plt.subplots(ncols=2, figsize=(10, 6),
                                           gridspec_kw=dict(width_ratios=[3, 1]))

        ax.plot(n_samples_arr, effect_size / means[index], '-ok')
        ax.plot(conv[index], eff_size_to_mean, marker='o', color="red")
        ax.text(conv[index] + 150, eff_size_to_mean + 0.01, '({}, {})'.format(conv[index], round(eff_size_to_mean, 5)))
        ax.set_title('Segment {}; mean = {}, std = {}, conv 30 days = {}'.format(segment_name, means[index], stds[index], conv[index]))
        ax.set_xlabel('Sample size')
        ax.set_ylabel('Effect size / Mean')

        ax_table.axis("off")
        cell_text = []
        for ind in reversed(range(len(n_samples_arr))):
            cell_text.append([round(n_samples_arr[ind]), round((effect_size / means[index])[ind], 5)])
        the_table = ax_table.table(cellText=cell_text,
                              colLabels=['Sample size', 'Effect size / Mean'],
                              loc='center left')

        plt.savefig('./output/plots/segment_' + str(segment) + '.png', dpi=300)
        plt.close()

    if create_presentation:
        # insert image to presentation
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = Inches(0)
        top = Inches(0)
        img = slide.shapes.add_picture('./output/plots/segment_' + str(segment) + '.png',
                                       left, top, width=Inches(16), height=Inches(9))

if create_presentation:
    prs.save('./presentation.pptx') # saving file

print('')
print('Check correctness of power test function:')
print(all(power_test))
